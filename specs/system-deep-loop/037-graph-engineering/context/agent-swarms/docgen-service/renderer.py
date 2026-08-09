"""
Native PowerPoint renderer (python-pptx) for AgentSwarms server-side doc-gen.

Takes a filled plan (the same PptxPlan the browser planner produces, with chart
categories/series already computed and each diagram slide carrying a pre-rendered
SVG) and produces a .pptx with NATIVE, EDITABLE charts, tables, KPI cards and
text — plus diagrams embedded as images (rasterised from the client's SVG).

The design MIRRORS the browser builder (src/lib/docGen/build.ts) slide for
slide: same geometry, type scale, cover treatment and chrome. That parity is the
point — "Deep" mode must never look worse than "Fast", because the only thing
the reader can compare is the deck. Where python-pptx has no direct equivalent
for a pptxgenjs feature (fill/text transparency, shrink-to-fit) we write the
DrawingML ourselves rather than dropping the effect.
"""
from __future__ import annotations

import io
import re
from datetime import date
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

CW = 13.333
CH = 7.5
M = 0.6
CONTENT_W = CW - M * 2
FONT = "Segoe UI"

INK = "0F172A"
INK_DEEP = "0B1220"
SUB = "64748B"
BODY = "334155"
BORDER = "E7EBF0"
CARD = "F8FAFC"
GRID = "EEF2F7"
DEFAULT_ACCENT = "4F46E5"
CHART_PALETTE = ["4F46E5", "0EA5E9", "10B981", "F59E0B", "EF4444", "8B5CF6", "EC4899", "14B8A6"]


def _hex(c: str | None, fallback: str = DEFAULT_ACCENT) -> str:
    h = re.sub(r"[^0-9a-fA-F]", "", (c or ""))
    return h.upper() if len(h) == 6 else fallback


def _rgb(c: str) -> RGBColor:
    return RGBColor.from_string(_hex(c))


def _lum(c: str) -> float:
    n = int(_hex(c), 16)
    return (0.2126 * ((n >> 16) & 255) + 0.7152 * ((n >> 8) & 255) + 0.0722 * (n & 255)) / 255


def _on_light(c: str) -> str:
    h = _hex(c)
    for _ in range(6):
        if _lum(h) <= 0.6:
            break
        n = int(h, 16)
        h = "".join(f"{round(((n >> s) & 255) * 0.72):02X}" for s in (16, 8, 0))
    return h


# ── DrawingML the python-pptx API doesn't expose ─────────────────────────────
def _set_fill_alpha(shape, opacity_pct: float) -> None:
    """Make a solid-filled shape translucent (pptxgenjs `transparency`)."""
    srgb = shape._element.spPr.find(qn("a:solidFill"))
    if srgb is None:
        return
    clr = srgb.find(qn("a:srgbClr"))
    if clr is None:
        return
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(opacity_pct * 1000))}))


def _set_text_alpha(run, opacity_pct: float) -> None:
    """Fade a run — used for the oversized section index watermark."""
    fill = run._r.get_or_add_rPr().find(qn("a:solidFill"))
    if fill is None:
        return
    clr = fill.find(qn("a:srgbClr"))
    if clr is None:
        return
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(opacity_pct * 1000))}))


def _autofit(text_frame) -> None:
    """
    PowerPoint's "shrink text on overflow". python-pptx's fit_text() needs font
    files and measures at render time; writing <a:normAutofit/> instead hands the
    job to the viewer, which is what pptxgenjs `fit: "shrink"` does. Without it
    an over-long title or bullet list silently runs off the slide.
    """
    bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        existing = bodyPr.find(qn(tag))
        if existing is not None:
            bodyPr.remove(existing)
    bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {}))


def _set_bullet(paragraph) -> None:
    """Square bullet with a hanging indent, so wrapped lines align to the text."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "228600")  # 0.25"
    pPr.set("indent", "-228600")
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial", "pitchFamily": "34", "charset": "0"}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "▪"}))


def _set_char_spacing(run, pts: float) -> None:
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def _is_number(v: Any) -> bool:
    if isinstance(v, (int, float)) and not isinstance(v, bool):
        return True
    return bool(re.fullmatch(r"[-+(]?[$£€]?\s?[\d,]+(\.\d+)?\s?[%)]?", str(v).strip()))


class Deck:
    def __init__(self, plan: dict[str, Any]):
        self.plan = plan
        self.accent = _hex(plan.get("accent"), DEFAULT_ACCENT)
        self.accent_ink = _on_light(self.accent)
        self.accent_dark = _shade(self.accent, 0.22)
        self.accent_light = _light(self.accent)
        self.title = plan.get("title") or "Untitled"
        self.date = date.today().strftime("%B %d, %Y")
        self.prs = Presentation()
        self.prs.slide_width = Inches(CW)
        self.prs.slide_height = Inches(CH)
        self.blank = self.prs.slide_layouts[6]
        self.slide_no = 0

    # ── low-level helpers ────────────────────────────────────────────────
    def _slide(self, bg: str | None = None):
        s = self.prs.slides.add_slide(self.blank)
        if bg:
            s.background.fill.solid()
            s.background.fill.fore_color.rgb = _rgb(bg)
        return s

    def _rect(self, slide, x, y, w, h, fill, *, radius=False, line=None, line_w=1.0, alpha=None,
              shape=MSO_SHAPE.RECTANGLE):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
        if alpha is not None:
            _set_fill_alpha(shp, alpha)
        if line:
            shp.line.color.rgb = _rgb(line)
            shp.line.width = Pt(line_w)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _text(self, slide, x, y, w, h, runs, *, size=14, color=BODY, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, inline=False,
              shrink=True, line_spacing=None):
        """
        `runs` is a string, or a list of (text, opts). A list normally means one
        paragraph per item (bullets); pass inline=True to keep them on ONE line
        — e.g. the "KEY INSIGHT <text>" bar, which otherwise wraps to two lines
        and overflows its 0.58" pill.
        """
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        if shrink:
            _autofit(tf)
        items = runs if isinstance(runs, list) else [(runs, {})]
        para = tf.paragraphs[0]
        first = True
        for text, opt in items:
            if not first and not inline:
                para = tf.add_paragraph()
            if first or not inline:
                para.alignment = opt.get("align", align)
                if opt.get("space_after") is not None:
                    para.space_after = Pt(opt["space_after"])
                if line_spacing:
                    para.line_spacing = line_spacing
                if opt.get("bullet"):
                    _set_bullet(para)
            first = False
            r = para.add_run()
            r.text = text
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.name = FONT
            r.font.color.rgb = _rgb(opt.get("color", color))
            if opt.get("alpha") is not None:
                _set_text_alpha(r, opt["alpha"])
            if spacing:
                _set_char_spacing(r, spacing)
        return tb

    def _svg_image(self, slide, svg: str, x, y, w, h, diagram=None):
        # Rasterise the diagram SVG. cairosvg needs system libs (cairo/pango);
        # if unavailable, fall back to the diagram's own text so the slide still
        # says something instead of showing an empty card.
        try:
            import cairosvg

            png = cairosvg.svg2png(
                bytestring=svg.encode("utf-8"),
                output_width=int(w * 150),
                output_height=int(h * 150),
            )
            slide.shapes.add_picture(io.BytesIO(png), Inches(x), Inches(y), Inches(w), Inches(h))
        except Exception:
            self._rect(slide, x, y, w, h, CARD, radius=True, line=BORDER)
            points = _diagram_text(diagram)
            if points:
                self._bullets(slide, points, x + 0.4, y + 0.35, w - 0.8, h - 0.7)

    # ── deck ─────────────────────────────────────────────────────────────
    def build(self) -> bytes:
        self._cover()
        self._contents()
        section_no = 0
        current_section = ""
        for s in self.plan.get("slides", []):
            layout = self._layout(s)
            if layout == "section":
                section_no += 1
                current_section = s.get("title") or f"Section {section_no}"
                self._section(section_no, s)
                continue
            slide = self._slide("FFFFFF")
            self._rect(slide, 0, 0, 0.16, CH, self.accent)  # accent spine
            self._header(slide, s.get("title", ""), current_section or self.title, s.get("subtitle"))
            self._footer(slide)
            bottom = 6.25 if s.get("takeaway") else 6.95
            top = 1.55
            self._content(slide, s, layout, top, bottom)
            if s.get("takeaway"):
                self._takeaway(slide, s["takeaway"])
            if s.get("notes"):
                slide.notes_slide.notes_text_frame.text = s["notes"]
        out = io.BytesIO()
        self.prs.save(out)
        return out.getvalue()

    def _layout(self, s: dict) -> str:
        if s.get("layout"):
            return s["layout"]
        if s.get("diagramSvg") or s.get("diagram"):
            return "diagram"
        if s.get("kpis"):
            return "kpi"
        if _chart_has_data(s.get("chart")):
            return "chart"
        if s.get("table"):
            return "table"
        return "bullets"

    def _cover(self):
        s = self._slide(INK_DEEP)
        # Concentric accent rings bleeding off the top-right corner for depth.
        self._rect(s, 8.9, -2.6, 7.4, 7.4, self.accent, alpha=20, shape=MSO_SHAPE.OVAL)
        self._rect(s, 10.3, -1.0, 5.2, 5.2, self.accent_dark, alpha=26, shape=MSO_SHAPE.OVAL)
        self._rect(s, 11.5, 3.3, 3.6, 3.6, "FFFFFF", alpha=8, shape=MSO_SHAPE.OVAL)
        self._rect(s, 0.9, 2.02, 0.62, 0.09, self.accent)
        self._text(s, 1.64, 1.82, 8, 0.4, f"REPORT · {self.date.upper()}", size=11.5,
                   bold=True, color=self.accent_light, spacing=3)
        self._text(s, 0.85, 2.45, 9.6, 2.3, self.title,
                   size=36 if len(self.title) > 48 else 46, bold=True, color="FFFFFF",
                   line_spacing=0.98)
        if self.plan.get("subtitle"):
            self._text(s, 0.9, 4.95, 8.6, 1.0, self.plan["subtitle"], size=18, color="CBD5E1")
        self._rect(s, 0.9, 6.45, 3.0, 0.02, "334155")
        n = len(self.plan.get("slides") or [])
        self._text(s, 0.9, 6.6, 8, 0.4, f"{n} slides · Generated {self.date}",
                   size=11, color="94A3B8")

    def _contents(self):
        """
        A contents page built from the deck's own section dividers. Only the
        server renderer produces this, so it is one of the things that visibly
        marks a Deep deck apart from the in-browser build — and on a 25-slide
        deck a reader genuinely needs it. Skipped when the deck is too small to
        warrant one.
        """
        sections = [
            s for s in self.plan.get("slides", [])
            if (s.get("layout") == "section") and s.get("title")
        ]
        if len(sections) < 3:
            return
        slide = self._slide("FFFFFF")
        self._rect(slide, 0, 0, 0.16, CH, self.accent)
        self._text(slide, M, 0.66, CONTENT_W, 0.7, "Contents", size=30, bold=True, color=INK)
        self._rect(slide, M, 1.42, 0.5, 0.05, self.accent)
        self._rect(slide, M + 0.6, 1.442, CONTENT_W - 0.6, 0.012, BORDER)

        # Two columns once a deck has more sections than fit comfortably.
        per_col = 6
        cols = 1 if len(sections) <= per_col else 2
        col_w = (CONTENT_W - (0.6 if cols == 2 else 0)) / cols
        rows = -(-len(sections) // cols)
        row_h = min(0.72, (6.6 - 1.9) / max(rows, 1))
        for i, sec in enumerate(sections[: per_col * 2]):
            c, r = divmod(i, rows)
            x = M + c * (col_w + 0.6)
            y = 1.9 + r * row_h
            self._text(slide, x, y, 0.6, row_h, f"{i + 1:02d}", size=17, bold=True,
                       color=self.accent_ink, anchor=MSO_ANCHOR.MIDDLE)
            self._text(slide, x + 0.62, y, col_w - 0.72, row_h, str(sec["title"]), size=15,
                       color=INK, anchor=MSO_ANCHOR.MIDDLE)
            self._rect(slide, x, y + row_h - 0.06, col_w - 0.1, 0.008, BORDER)
        self._footer(slide)

    def _section(self, no: int, s: dict):
        slide = self._slide(INK_DEEP)
        # Oversized, FAINT index bleeding off the bottom-right. Without the alpha
        # this reads as a solid block of colour that dominates the slide.
        self._text(slide, 6.4, 1.1, 6.6, 6.2, [(f"{no:02d}", {"alpha": 14})], size=300, bold=True,
                   color=self.accent_light, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM,
                   shrink=False)
        self._rect(slide, 0.9, 2.72, 0.62, 0.09, self.accent_light)
        title = s.get("title", "")
        self._text(slide, 1.64, 2.52, 8, 0.4, f"SECTION {no:02d}", size=11.5, bold=True,
                   color=self.accent_light, spacing=3)
        self._text(slide, 0.9, 3.02, 9.4, 1.5, title, size=30 if len(title) > 42 else 38,
                   bold=True, color="FFFFFF")
        if s.get("subtitle"):
            self._text(slide, 0.92, 4.6, 8.8, 0.8, s["subtitle"], size=15, color="CBD5E1")
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]

    def _header(self, slide, title: str, kicker: str, subtitle: str | None = None):
        if kicker:
            self._text(slide, M, 0.4, CONTENT_W, 0.26, kicker.upper()[:60], size=10.5,
                       bold=True, color=self.accent_ink, spacing=2.5)
        title = title or ""
        # Two lines of 24pt overrun the rule at 1.32"; step down like the browser
        # builder does, and let autofit handle anything still too long.
        self._text(slide, M - 0.02, 0.64 if kicker else 0.5,
                   CONTENT_W - 3.6 if subtitle else CONTENT_W, 0.62,
                   title, size=20 if len(title) > 60 else 24, bold=True, color=INK)
        if subtitle:
            self._text(slide, M + CONTENT_W - 3.5, 0.78 if kicker else 0.64, 3.5, 0.34,
                       subtitle, size=11.5, color=SUB, align=PP_ALIGN.RIGHT)
        self._rect(slide, M, 1.3, 0.5, 0.05, self.accent)
        self._rect(slide, M + 0.6, 1.322, CONTENT_W - 0.6, 0.012, BORDER)

    def _footer(self, slide):
        self.slide_no += 1
        self._text(slide, 0.5, 7.08, 9, 0.32, self.title, size=8, color=SUB,
                   anchor=MSO_ANCHOR.MIDDLE, shrink=False)
        self._text(slide, 12.4, 7.08, 0.6, 0.32, str(self.slide_no), size=8, color=SUB,
                   align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, shrink=False)

    def _takeaway(self, slide, text: str):
        self._rect(slide, M, 6.4, CONTENT_W, 0.58, _tint(self.accent, 0.92), radius=True)
        self._rect(slide, M, 6.4, 0.11, 0.58, self.accent, radius=True)
        self._text(slide, M + 0.32, 6.4, CONTENT_W - 0.6, 0.58,
                   [("KEY INSIGHT    ", {"bold": True, "color": self.accent_ink}),
                    (str(text)[:240], {"color": INK})],
                   size=12, anchor=MSO_ANCHOR.MIDDLE, inline=True)

    # ── content dispatch ─────────────────────────────────────────────────
    def _content(self, slide, s, layout, top, bottom):
        if layout == "diagram" and s.get("diagramSvg"):
            self._svg_image(slide, s["diagramSvg"], M, top, CONTENT_W, bottom - top,
                            diagram=s.get("diagram"))
            return
        if layout == "kpi" and s.get("kpis"):
            self._kpis(slide, s["kpis"][:5], s.get("bullets"), bottom)
            return
        has_chart = _chart_has_data(s.get("chart"))
        has_table = _table_has_data(s.get("table"))
        has_visual = has_chart or has_table
        has_text = bool(s.get("bullets") or s.get("paragraph"))
        if not has_visual and not has_text:
            # A "chart" slide whose query returned nothing, or one the verify
            # loop stripped: the content area would be white space under a
            # title. Never ship that — say what the slide is about instead.
            self._empty_state(slide, s, top, bottom)
            return
        if has_visual and has_text:
            if has_chart:
                self._chart(slide, s["chart"], M, top, 7.3, bottom - top)
            else:
                self._table(slide, s["table"], M, top, 7.3, bottom - top, 12)
            ty = top + 0.05
            if s.get("paragraph"):
                self._text(slide, 8.2, ty, 4.5, 1.6, s["paragraph"], size=14, color=BODY)
                ty += 1.7
            if s.get("bullets"):
                self._bullets(slide, s["bullets"], 8.2, ty, 4.5, bottom - ty)
        elif has_visual:
            if has_chart:
                self._chart(slide, s["chart"], M, top, CONTENT_W, bottom - top)
            else:
                self._table(slide, s["table"], M, top, CONTENT_W, bottom - top, 13)
        else:
            y = top
            if s.get("paragraph"):
                self._text(slide, M, y, CONTENT_W, 1.6, s["paragraph"], size=16, color=BODY)
                y += 1.7
            if s.get("bullets"):
                self._bullets(slide, s["bullets"], M, y, CONTENT_W, bottom - y)

    def _bullets(self, slide, bullets, x, y, w, h):
        runs = [(str(b), {"bullet": True, "size": 18, "color": BODY, "space_after": 15})
                for b in bullets]
        self._text(slide, x, y, w, h, runs)

    def _empty_state(self, slide, s, top, bottom):
        """Last-resort content so no slide ever renders as a bare title."""
        notes = str(s.get("notes") or "").strip()
        takeaway = str(s.get("takeaway") or "").strip()
        body = notes or takeaway
        if body:
            self._text(slide, M, top + 0.1, CONTENT_W, min(2.4, bottom - top - 0.2),
                       body, size=16, color=BODY)
            return
        h = min(1.1, bottom - top)
        self._rect(slide, M, top + 0.1, CONTENT_W, h, CARD, radius=True, line=BORDER)
        self._text(slide, M + 0.4, top + 0.1, CONTENT_W - 0.8, h,
                   "No data was available for this view.", size=13, color=SUB,
                   anchor=MSO_ANCHOR.MIDDLE)

    def _card(self, slide, x, y, w, h, top_rule=False):
        self._rect(slide, x, y, w, h, "FFFFFF", radius=True, line=BORDER)
        if top_rule:
            self._rect(slide, x, y, w, 0.09, self.accent)

    def _kpis(self, slide, kpis, bullets, bottom):
        gap = 0.28
        cw = (CONTENT_W - gap * (len(kpis) - 1)) / len(kpis)
        cy, ch = 1.95, 2.3
        for i, k in enumerate(kpis):
            x = M + i * (cw + gap)
            self._card(slide, x, cy, cw, ch, top_rule=True)
            self._text(slide, x + 0.26, cy + 0.34, cw - 0.5, 0.7, str(k.get("label", "")).upper(),
                       size=10.5, bold=True, color=SUB, spacing=1.2)
            self._text(slide, x + 0.24, cy + 0.9, cw - 0.42, 0.9, str(k.get("value", "")),
                       size=33, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
            if k.get("delta"):
                good = k.get("positive") is not False
                delta = str(k["delta"])
                # Pill sized to its text — a fixed width wraps "+12.4% vs FY25"
                # onto two lines and spills outside the pill.
                pw = min(cw - 0.5, 0.42 + len(delta) * 0.11)
                self._rect(slide, x + 0.26, cy + ch - 0.62, pw, 0.34,
                           "DCFCE7" if good else "FEE2E2", radius=True)
                self._text(slide, x + 0.26, cy + ch - 0.62, pw, 0.34, delta,
                           size=10.5, bold=True, color="047857" if good else "B91C1C",
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if bullets:
            self._bullets(slide, bullets, M, 4.55, CONTENT_W, bottom - 4.55)

    def _table(self, slide, t, x, y, w, max_h, font_pt):
        cols = t.get("columns", [])
        rows = t.get("rows", [])
        if not cols:
            return
        # Keep the table inside the content area: a plan with more rows than fit
        # would otherwise run off the slide and under the key-insight bar.
        row_h = 0.34 if font_pt >= 13 else 0.3
        max_rows = max(1, int((max_h - row_h) / row_h))
        clipped = len(rows) > max_rows
        shown = rows[: max_rows - 1] if clipped else rows
        n = len(shown) + 1 + (1 if clipped else 0)
        gshape = slide.shapes.add_table(n, len(cols), Inches(x), Inches(y), Inches(w),
                                        Inches(min(row_h * n, max_h)))
        table = gshape.table
        for c, col in enumerate(cols):
            cell = table.cell(0, c)
            cell.text = str(col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(self.accent)
            for pr in cell.text_frame.paragraphs:
                for rn in pr.runs:
                    rn.font.bold = True
                    rn.font.color.rgb = _rgb("FFFFFF")
                    rn.font.size = Pt(font_pt)
                    rn.font.name = FONT
        for r, row in enumerate(shown):
            for c in range(len(cols)):
                val = row[c] if c < len(row) else ""
                cell = table.cell(r + 1, c)
                cell.text = "" if val is None else str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("FFFFFF" if r % 2 == 0 else CARD)
                for pr in cell.text_frame.paragraphs:
                    if c > 0 and _is_number(val):
                        pr.alignment = PP_ALIGN.RIGHT
                    for rn in pr.runs:
                        rn.font.size = Pt(font_pt - 1)
                        rn.font.name = FONT
                        rn.font.color.rgb = _rgb(BODY)
        if clipped:
            cell = table.cell(n - 1, 0)
            cell.text = f"… {len(rows) - len(shown)} more rows"
            for pr in cell.text_frame.paragraphs:
                for rn in pr.runs:
                    rn.font.size = Pt(font_pt - 1)
                    rn.font.name = FONT
                    rn.font.italic = True
                    rn.font.color.rgb = _rgb(SUB)

    def _chart(self, slide, chart, x, y, w, h):
        self._card(slide, x, y, w, h)
        cd = CategoryChartData()
        cd.categories = chart.get("categories", [])
        series = chart.get("series", []) or []
        for ser in series:
            cd.add_series(ser.get("name", "Series"), tuple(float(v or 0) for v in ser.get("values", [])))
        kind = chart.get("type")
        ctype = {
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "area": XL_CHART_TYPE.AREA,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        }.get(kind, XL_CHART_TYPE.COLUMN_CLUSTERED)
        gf = slide.shapes.add_chart(ctype, Inches(x + 0.25), Inches(y + 0.25),
                                    Inches(w - 0.5), Inches(h - 0.5), cd)
        ch = gf.chart
        ch.has_title = False
        # Chart-wide type scale. Left at PowerPoint's defaults the axis labels and
        # legend come out oversized in a wrong font — the "stock Excel chart" look
        # the browser builder avoids by drawing its own SVG.
        ch.font.size = Pt(10)
        ch.font.name = FONT
        ch.font.color.rgb = _rgb(SUB)
        is_pie = kind in ("pie", "doughnut")
        ch.has_legend = is_pie or len(series) > 1
        if ch.has_legend:
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
            ch.legend.font.size = Pt(10)
            ch.legend.font.name = FONT
            ch.legend.font.color.rgb = _rgb(BODY)
        if not is_pie:
            self._style_axes(ch)
        is_line = kind in ("line", "area")
        # One colour per CATEGORY only where that reads as data (a pie, or a
        # single-series bar); a single-series line must stay one colour.
        vary = is_pie or (len(series) == 1 and not is_line)
        try:
            plot = ch.plots[0]
            plot.vary_by_categories = vary
            if kind in (None, "bar", "column") or ctype in (
                XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.BAR_CLUSTERED
            ):
                plot.gap_width = 60
                if len(series) > 1:
                    plot.overlap = -10
            if is_pie:
                plot.has_data_labels = True
                labels = plot.data_labels
                labels.show_percentage = True
                labels.show_value = False
                labels.number_format = "0.0%"
                labels.number_format_is_linked = False
                labels.font.size = Pt(10)
                labels.font.name = FONT
                labels.font.color.rgb = _rgb("FFFFFF")
            for idx, ser in enumerate(plot.series):
                colour = _rgb(CHART_PALETTE[idx % len(CHART_PALETTE)])
                if vary:
                    # With varyColors set, PowerPoint ignores the SERIES fill and
                    # colours each point from the theme — which is how a pie ends
                    # up in the stock blue/red/green palette. Colour the points.
                    for pi, point in enumerate(ser.points):
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = _rgb(
                            CHART_PALETTE[pi % len(CHART_PALETTE)]
                        )
                        point.format.line.fill.background()
                elif is_line:
                    ser.format.line.color.rgb = colour
                    ser.format.line.width = Pt(2.25)
                    ser.smooth = False
                else:
                    ser.format.fill.solid()
                    ser.format.fill.fore_color.rgb = colour
                    ser.format.line.fill.background()
        except Exception:
            pass

    def _style_axes(self, ch) -> None:
        try:
            va = ch.value_axis
            va.has_major_gridlines = True
            va.major_gridlines.format.line.color.rgb = _rgb(GRID)
            va.major_gridlines.format.line.width = Pt(0.75)
            va.format.line.fill.background()
            va.has_title = False
            va.tick_labels.font.size = Pt(10)
            va.tick_labels.font.name = FONT
            va.tick_labels.font.color.rgb = _rgb(SUB)
            # Thousands separators — raw "100000" tick labels are the clearest
            # tell of an unstyled chart.
            va.tick_labels.number_format = "#,##0"
            va.tick_labels.number_format_is_linked = False
        except Exception:
            pass
        try:
            ca = ch.category_axis
            ca.has_major_gridlines = False
            ca.has_title = False
            ca.format.line.color.rgb = _rgb(BORDER)
            ca.tick_labels.font.size = Pt(10)
            ca.tick_labels.font.name = FONT
            ca.tick_labels.font.color.rgb = _rgb(SUB)
        except Exception:
            pass


def _light(c: str) -> str:
    """Lighten toward white for readable text on the dark cover/section bg."""
    h = _hex(c)
    for _ in range(6):
        if _lum(h) >= 0.62:
            break
        n = int(h, 16)
        h = "".join(f"{round(v + (255 - v) * 0.35):02X}" for v in ((n >> 16) & 255, (n >> 8) & 255, n & 255))
    return h


def _shade(c: str, amt: float) -> str:
    n = int(_hex(c), 16)
    return "".join(f"{round(v * (1 - amt)):02X}" for v in ((n >> 16) & 255, (n >> 8) & 255, n & 255))


def _tint(c: str, amt: float) -> str:
    n = int(_hex(c), 16)
    return "".join(
        f"{round(v + (255 - v) * amt):02X}" for v in ((n >> 16) & 255, (n >> 8) & 255, n & 255)
    )


def _chart_has_data(chart) -> bool:
    if not chart:
        return False
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    return len(cats) > 0 and any(len(s.get("values") or []) > 0 for s in series)


def _table_has_data(table) -> bool:
    """A table counts as a visual only if it has columns AND at least one row.

    When a chart's query returns nothing the client substitutes a fallback
    table, and that substitute can itself come back empty. Treating it as a
    visual regardless is what put an empty grid next to the bullets instead of
    reflowing the slide or showing the empty state.
    """
    if not isinstance(table, dict):
        return False
    cols = table.get("columns") or []
    rows = table.get("rows") or []
    return len(cols) > 0 and len(rows) > 0


def _diagram_text(diagram) -> list[str]:
    """Headline strings from a diagram spec, for the no-rasteriser fallback."""
    if not isinstance(diagram, dict):
        return []
    out: list[str] = []
    for key in ("steps", "cards", "stages", "tiers", "quadrants", "phases", "nodes", "sets"):
        for item in diagram.get(key) or []:
            if isinstance(item, dict):
                label = item.get("title") or item.get("label") or ""
                detail = item.get("detail")
                out.append(f"{label} — {detail}" if label and detail else str(label))
    for col in diagram.get("columns") or []:
        if isinstance(col, dict):
            out.append(str(col.get("heading") or col.get("title") or ""))
    return [t for t in out if t][:8]


def render_pptx(plan: dict[str, Any]) -> bytes:
    return Deck(plan).build()
